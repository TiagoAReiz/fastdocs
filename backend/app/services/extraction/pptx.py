import io

from pptx import Presentation

from app.services.extraction.base import ExtractionResult, Section


def extract(content: bytes) -> ExtractionResult:
    prs = Presentation(io.BytesIO(content))
    sections: list[Section] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        text = "\n".join(parts).strip()
        if text:
            sections.append(
                Section(text=text, metadata={"slide_number": i}, kind="text")
            )
    return ExtractionResult(
        sections=sections,
        extraction_method="native",
        document_metadata={"slide_count": len(prs.slides)},
    )
